from pptx import Presentation
from .models import SlidePlan
import random, tempfile

def generate_pptx(template_file, slide_plan: SlidePlan) -> str:
    prs = Presentation(template_file)

    # Collect layouts
    layouts = {l.name.lower(): l for l in prs.slide_layouts}

    # Collect template images
    template_images = [p for p in prs.part.related_parts.values()
                       if getattr(p, "content_type", "").startswith("image/")]

    for slide_content in slide_plan.slides:
        layout = None
        for key in layouts:
            if slide_content.layout_hint in key:
                layout = layouts[key]
                break
        if layout is None:
            layout = prs.slide_layouts[1]  # fallback title+content

        slide = prs.slides.add_slide(layout)

        # Fill title
        if slide.shapes.title:
            slide.shapes.title.text = slide_content.title

        # Fill bullets
        body_shapes = [s for s in slide.placeholders if s.placeholder_format.type == 1]
        if body_shapes:
            tx = body_shapes[0].text_frame
            tx.clear()
            for i, bullet in enumerate(slide_content.bullets[:slide_plan.max_bullets_per_slide]):
                p = tx.add_paragraph() if i else tx.paragraphs[0]
                p.text = " ".join(bullet.split()[:slide_plan.max_words_per_bullet])
                p.level = 0

        # Optional: insert a template image randomly
        picture_placeholders = [s for s in slide.placeholders if "picture" in s.name.lower()]
        if picture_placeholders and template_images:
            img = random.choice(template_images)
            with tempfile.NamedTemporaryFile(delete=False) as tmp_img:
                tmp_img.write(img.blob)
                picture_placeholders[0].insert_picture(tmp_img.name)

        # Speaker notes
        if slide.has_notes_slide and slide_content.notes:
            slide.notes_slide.notes_text_frame.text = slide_content.notes

    output = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(output.name)
    return output.name
